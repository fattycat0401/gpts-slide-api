from pptx import Presentation
from pptx.util import Inches, Pt
from flask import Flask, request, send_from_directory, jsonify
from threading import Thread
import os
import json
import uuid
import time

app = Flask(__name__)

### 設定參數 ###
TOKEN = "fattycat0401"
STATIC_FOLDER = "static"
EXPIRATION_TIME = 600  # 10分鐘
os.makedirs(STATIC_FOLDER, exist_ok=True)


### 檔案清理背景程序 ###
def cleanup_expired_files_safe():
    while True:
        now = time.time()
        for filename in os.listdir(STATIC_FOLDER):
            filepath = os.path.join(STATIC_FOLDER, filename)
            if os.path.isfile(filepath) and now - os.path.getmtime(filepath) > EXPIRATION_TIME:
                try:
                    os.remove(filepath)
                except Exception:
                    pass
        time.sleep(60)  # 每分鐘掃描一次


Thread(target=cleanup_expired_files_safe, daemon=True).start()


### 工具函式：建立首頁 ###
def add_title_slide(prs, title_text):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(2))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title_text
    run.font.size = Pt(48)
    run.font.bold = True
    return slide


### 工具函式：建立內容或結尾頁 ###
def add_content_slide(prs, h2, sections):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    top = Inches(0.5)

    # H2 標題
    h2_box = slide.shapes.add_textbox(Inches(0.5), top, Inches(9), Inches(1))
    h2_frame = h2_box.text_frame
    h2_frame.word_wrap = True
    h2_run = h2_frame.paragraphs[0].add_run()
    h2_run.text = h2
    h2_run.font.size = Pt(36)
    h2_run.font.bold = True
    top += Inches(1.0)

    for section in sections:
        if section.get("h3"):
            h3_box = slide.shapes.add_textbox(Inches(0.7), top, Inches(8), Inches(0.5))
            h3_frame = h3_box.text_frame
            h3_frame.word_wrap = True
            h3_run = h3_frame.paragraphs[0].add_run()
            h3_run.text = section["h3"]
            h3_run.font.size = Pt(24)
            h3_run.font.bold = True
            top += Inches(0.5)

        if section.get("p"):
            p_box = slide.shapes.add_textbox(Inches(1.0), top, Inches(8.5), Inches(1.2))
            p_frame = p_box.text_frame
            p_frame.word_wrap = True
            p_run = p_frame.paragraphs[0].add_run()
            p_run.text = section["p"]
            p_run.font.size = Pt(18)
            top += Inches(1.0)

    return slide


### 主頁面測試用 ###
@app.route("/")
def home():
    return "PPT Generator API is running."


### 產生簡報主邏輯 ###
@app.route("/generate_pptx", methods=["POST"])
def generate_pptx():
    data = request.get_json()
    token = data.get("token", "")
    if token != TOKEN:
        return jsonify({"error": "Invalid token"}), 403

    h1 = data.get("h1", "Untitled Presentation")
    pages = data.get("pages", [])

    prs = Presentation()

    # 首頁
    add_title_slide(prs, h1)

    # 內容與結尾頁
    for page in pages:
        h2 = page.get("h2", "")
        sections = page.get("sections", [])
        add_content_slide(prs, h2, sections)

    # 儲存檔案
    filename = f"{uuid.uuid4().hex}.pptx"
    filepath = os.path.join(STATIC_FOLDER, filename)
    prs.save(filepath)

    return jsonify({"download_url": f"/static/{filename}"})


### 提供下載簡報 ###
@app.route("/static/<filename>")
def download_file(filename):
    return send_from_directory(STATIC_FOLDER, filename, as_attachment=True)
